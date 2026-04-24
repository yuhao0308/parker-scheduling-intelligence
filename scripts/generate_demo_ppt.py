from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("docs/united_hebrew_demo_minimalist_presentation.pptx")

BG = RGBColor(246, 244, 239)
SURFACE = RGBColor(255, 255, 255)
TEXT = RGBColor(27, 32, 37)
MUTED = RGBColor(103, 110, 118)
ACCENT = RGBColor(27, 86, 112)
ACCENT_SOFT = RGBColor(228, 238, 243)
SUCCESS_SOFT = RGBColor(232, 241, 236)
WARN_SOFT = RGBColor(245, 236, 227)
LINE = RGBColor(215, 219, 222)


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.72), Inches(0.48), Inches(11.4), Inches(0.95))
    tf = box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = TEXT

    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(11)
        r2.font.color.rgb = MUTED

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(1.42), Inches(1.45), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_card(slide, left, top, width, height, title, lines, fill_rgb=SURFACE) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = LINE

    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(9)
    tf.margin_bottom = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.bold = True
    r.font.size = Pt(15)
    r.font.color.rgb = TEXT
    p.space_after = Pt(5)

    for item in lines:
        p = tf.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT
        p.space_after = Pt(1.5)


def add_formula_card(slide, left, top, width, height, title, formula_lines, footnote=None) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_SOFT
    shape.line.color.rgb = LINE

    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = TEXT
    p.space_after = Pt(7)

    for line in formula_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos Mono"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT
        p.space_after = Pt(1.5)

    if footnote:
        p = tf.add_paragraph()
        p.text = footnote
        p.font.name = "Aptos"
        p.font.size = Pt(10)
        p.font.color.rgb = MUTED


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.78), Inches(6.88), Inches(11.0), Inches(0.28))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def add_notes(slide, notes_text: str) -> None:
    notes_slide = slide.notes_slide
    for placeholder in notes_slide.placeholders:
        if placeholder.placeholder_format.idx == 3:
            tf = placeholder.text_frame
            tf.clear()
            paragraphs = notes_text.strip().split("\n")
            first = True
            for line in paragraphs:
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.name = "Aptos"
                p.font.size = Pt(12)
            break


def slide_1(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "United Hebrew Scheduling Intelligence", "Scheduler math and recommendation math, presented as an explainable intelligence layer")
    add_card(
        slide,
        Inches(0.82), Inches(1.78), Inches(3.7), Inches(3.05),
        "Business focus",
        [
            "Monthly schedule generation",
            "Same-day call-out replacement",
            "Lower avoidable overtime",
            "More consistent staffing decisions",
        ],
    )
    add_card(
        slide,
        Inches(4.82), Inches(1.78), Inches(3.55), Inches(3.05),
        "What makes it credible",
        [
            "Rule-driven filtering before scoring",
            "Weighted ranking with visible math",
            "Plain-language rationale for each output",
            "Human decision remains in control",
        ],
        fill_rgb=ACCENT_SOFT,
    )
    add_card(
        slide,
        Inches(8.67), Inches(1.78), Inches(3.75), Inches(3.05),
        "Anchor employee used in both examples",
        [
            "Maria Santos (CNA013)",
            "Full-Time",
            "Home unit: U-LT4",
            "Home typology: Long-Term",
            "Hire date: 2018-10-01",
        ],
        fill_rgb=SUCCESS_SOFT,
    )
    add_footer(slide, "Every score component is normalized to 0-1 before applying weights.")
    add_notes(
        slide,
        """
This demo is designed to show two things clearly.
First, the monthly scheduler is not assigning people randomly. It is using explicit rules and explicit weighted math.
Second, the recommendation engine uses the same explainable logic, but with more real-time signals for urgent call-out decisions.

To make the math concrete, I will use one seeded employee example in both workflows: Maria Santos, employee ID CNA013.
She is a full-time CNA whose home unit is U-LT4, which is a Long-Term unit, and her hire date is October 1, 2018.

That way, the audience does not just hear the factor names. They can see how one actual employee would be evaluated and how the score is calculated step by step.
        """,
    )


def slide_2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "System View", "Two workflows, one explainable scoring foundation")
    add_card(
        slide,
        Inches(0.82), Inches(1.75), Inches(3.7), Inches(4.65),
        "Monthly scheduler",
        [
            "Input: staff pool, units, hours ledger, home unit, cross-training, scoring weights",
            "Decision path: filter -> score -> assign",
            "Primary signals: overtime headroom, clinical fit, float penalty",
            "Output: best possible month from the current database state",
        ],
    )
    add_card(
        slide,
        Inches(4.82), Inches(1.75), Inches(3.55), Inches(4.65),
        "Recommendation engine",
        [
            "Input: call-out employee, target shift, PTO, exclusions, current schedule context",
            "Decision path: filter -> score -> rank -> explain",
            "Primary signals: scheduler signals plus proximity, seniority, equity, willingness",
            "Output: ranked replacement list with rationale",
        ],
        fill_rgb=ACCENT_SOFT,
    )
    add_formula_card(
        slide,
        Inches(8.67), Inches(1.75), Inches(3.75), Inches(4.65),
        "Shared philosophy",
        [
            "1. Remove unsafe options first",
            "2. Score only the eligible staff",
            "3. Rank by weighted total",
            "4. Keep the final decision human",
        ],
        "The purpose of the demo is not to present a black box. It is to present an explainable staffing engine."
    )
    add_footer(slide, "This structure is useful in demos because it matches how customers naturally ask questions.")
    add_notes(
        slide,
        """
Before going into the formulas, I want to frame the system at a high level.
There are two workflows.

The monthly scheduler is for planned staffing. It starts with the month, the units, and the available staff, and it decides who should be assigned shift by shift.

The recommendation engine is for urgent staffing. It starts when someone calls out, then it removes everyone who is not eligible, scores the valid people, ranks them, and explains why they appear in that order.

The important point is that both workflows share the same philosophy:
first filter, then score, then rank, while keeping the final decision human-controlled.
        """,
    )


def slide_3(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Monthly Scheduler", "Decision rules and exact score formula")
    add_card(
        slide,
        Inches(0.82), Inches(1.74), Inches(3.7), Inches(4.75),
        "Eligibility before scoring",
        [
            "Preferred bucket by unit type",
            "Subacute -> licensed first",
            "Long-Term -> certified first",
            "Then remove: wrong bucket, already scheduled, rest-window conflict, weekly-cap violation",
        ],
    )
    add_formula_card(
        slide,
        Inches(4.82), Inches(1.74), Inches(3.55), Inches(2.35),
        "Weekly cap test",
        [
            "current_weekly_hours + 8.25 <= weekly_cap",
            "FT cap = 41.25",
            "PT cap = 24.75",
            "PER_DIEM cap = 41.25",
        ],
    )
    add_formula_card(
        slide,
        Inches(4.82), Inches(4.25), Inches(3.55), Inches(2.25),
        "Monthly score",
        [
            "MonthlyScore = 0.45*OT",
            "              + 0.22*ClinicalFit",
            "              - 0.13*FloatPenalty",
            "              + 0.05*0.5",
        ],
        "Proximity is fixed to 0.5 in monthly scheduling, so it is effectively neutral for ranking."
    )
    add_card(
        slide,
        Inches(8.67), Inches(1.74), Inches(3.75), Inches(4.75),
        "Interpretation",
        [
            "OT headroom carries the largest weight.",
            "Clinical fit protects the unit context.",
            "Float penalty discourages unnecessary movement.",
            "If nobody survives the filters, the gap stays visible rather than being hidden by a bad assignment.",
        ],
        fill_rgb=SUCCESS_SOFT,
    )
    add_footer(slide, "Operational day runs 11 PM to 11 PM, so Night shift belongs to the next operational day.")
    add_notes(
        slide,
        """
This slide shows the exact monthly scheduling logic.

The scheduler does not start by scoring everybody.
It first determines the preferred staffing bucket for the unit type.
For Subacute it prefers licensed staff first.
For Long-Term it prefers certified staff first.

Then it applies hard filters.
If a person is in the wrong bucket, already scheduled, blocked by the rest window, or would break the weekly cap, they are removed completely.

Only after that do we score the surviving candidates.
The monthly formula is:
MonthlyScore equals 0.45 times overtime headroom, plus 0.22 times clinical fit, minus 0.13 times float penalty, plus 0.05 times 0.5.

That last term comes from proximity, but in monthly scheduling proximity is fixed to 0.5 for everyone, so it is effectively neutral.
In practical terms, the monthly scheduler is mainly driven by overtime headroom, clinical fit, and float penalty.
        """,
    )


def slide_4(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Monthly Scheduler Example", "Maria Santos scored for one Long-Term shift")
    add_card(
        slide,
        Inches(0.82), Inches(1.75), Inches(3.55), Inches(4.65),
        "Scenario",
        [
            "Target slot: U-LT4 Day",
            "Unit type: Long-Term",
            "Maria license: CNA -> correct preferred bucket",
            "Generated-week hours so far: 8.25",
            "Hours ledger this cycle: 24.75",
            "Home unit matches target unit",
        ],
    )
    add_formula_card(
        slide,
        Inches(4.62), Inches(1.75), Inches(3.35), Inches(4.65),
        "Maria's component values",
        [
            "Weekly cap check:",
            "8.25 + 8.25 <= 41.25 -> pass",
            "",
            "hours_for_scoring = 24.75 + 8.25 = 33.0",
            "OT = 1.0",
            "ClinicalFit = 1.0",
            "FloatPenalty = 0.0",
            "Proximity = 0.5",
        ],
    )
    add_formula_card(
        slide,
        Inches(8.27), Inches(1.75), Inches(4.15), Inches(4.65),
        "Maria's final score",
        [
            "MonthlyScore = 0.45(1.0)",
            "              + 0.22(1.0)",
            "              - 0.13(0.0)",
            "              + 0.05(0.5)",
            "",
            "            = 0.45 + 0.22 + 0 + 0.025",
            "            = 0.695",
            "",
            "Displayed intuition: 69.5/100",
        ],
        "Maria scores well because she still has straight-time headroom, she is a perfect unit fit, and she creates no float penalty."
    )
    add_footer(slide, "The scheduler repeats this calculation for every surviving candidate and assigns the highest total.")
    add_notes(
        slide,
        """
Now I want to show the actual math using one employee.

Suppose we are filling a Day shift in U-LT4, a Long-Term unit.
Maria Santos is a CNA, so she is in the correct preferred bucket for Long-Term.

Before we even score her, we test the weekly cap.
She already has 8.25 generated hours in that week, and adding one more shift makes 16.5.
That is still below the full-time cap of 41.25, so she passes.

For scoring, her cycle hours plus current generated-week hours give 33.0.
Because 33.0 is still below the 37.5 weekly overtime threshold, her overtime headroom is 1.0.
Her clinical fit is 1.0 because U-LT4 is her home unit.
Her float penalty is 0 because she is not floating away from home.
Proximity is fixed to 0.5 in monthly scheduling.

So the final math is:
0.45 times 1.0, plus 0.22 times 1.0, minus 0.13 times 0, plus 0.05 times 0.5.
That equals 0.695.

This is the exact kind of walkthrough that helps the audience understand how the scheduler chooses one person over another.
        """,
    )


def slide_5(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Recommendation Engine", "Real-time filters and full ranking formula")
    add_card(
        slide,
        Inches(0.82), Inches(1.74), Inches(3.7), Inches(4.8),
        "Hard filters first",
        [
            "Infer licensed or certified bucket from the employee who called out",
            "Remove wrong bucket",
            "Remove already scheduled staff",
            "Remove PTO conflicts and unit exclusions",
            "Remove rest-window conflicts",
            "Remove candidates whose removal would break source-unit minimum coverage",
        ],
    )
    add_formula_card(
        slide,
        Inches(4.82), Inches(1.74), Inches(3.55), Inches(4.8),
        "Recommendation score",
        [
            "RecScore = 0.45*OT",
            "         + 0.22*ClinicalFit",
            "         + 0.05*Proximity",
            "         + 0.08*Seniority",
            "         + 0.05*Equity",
            "         + 0.02*Willingness",
            "         - 0.13*FloatPenalty",
        ],
        "The UI displays round(RecScore * 100)."
    )
    add_card(
        slide,
        Inches(8.67), Inches(1.74), Inches(3.75), Inches(4.8),
        "Signal meaning",
        [
            "OT: avoid avoidable overtime",
            "Clinical fit: reward home unit, cross-training, and safe typology match",
            "Proximity: favor faster response, but only lightly",
            "Seniority, equity, willingness: refine rather than dominate",
        ],
        fill_rgb=WARN_SOFT,
    )
    add_footer(slide, "The ranking is transparent because each score is broken into named components.")
    add_notes(
        slide,
        """
The recommendation engine uses the same idea, but it is more real-time.

When someone calls out, the first step is to infer the replacement bucket from that employee.
If the call-out is certified, we only consider certified staff.
If the call-out is licensed, we only consider licensed staff.

Then we apply hard filters again, but now with live context:
already scheduled staff, PTO conflicts, unit exclusions, rest-window conflicts, and source-unit coverage floors.

Only then do we apply the full weighted score.
The recommendation formula includes the monthly signals plus more real-time refinement:
proximity, seniority, equity, and willingness.

The important message to say aloud is that overtime is still the dominant factor at 45 percent, and clinical fit is still next at 22 percent.
The smaller factors improve the ranking quality without overpowering the core staffing priorities.
        """,
    )


def slide_6(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Recommendation Example", "Maria Santos scored for a certified Long-Term call-out")
    add_card(
        slide,
        Inches(0.82), Inches(1.72), Inches(3.45), Inches(4.9),
        "Scenario",
        [
            "Call-out: U-LT1 Day on 2026-04-17",
            "Required bucket: certified",
            "Maria is eligible for the bucket",
            "No PTO conflict",
            "No unit exclusion",
            "Distance from 11432 to facility ZIP 11375: 2.87 miles",
        ],
    )
    add_formula_card(
        slide,
        Inches(4.47), Inches(1.72), Inches(3.75), Inches(4.9),
        "Maria's normalized values",
        [
            "OT = 1.0",
            "ClinicalFit = 1.0",
            "Proximity = 1 - 2.87/30 = 0.9043",
            "FloatPenalty = 0.3",
            "Seniority = 7.54/10 = 0.754",
            "Equity = 0.0",
            "Willingness = 0.5",
        ],
        "Float penalty is 0.3 because the target unit differs from home unit but remains in the same Long-Term typology."
    )
    add_formula_card(
        slide,
        Inches(8.47), Inches(1.72), Inches(3.95), Inches(4.9),
        "Maria's final score",
        [
            "RecScore = 0.45(1.0)",
            "         + 0.22(1.0)",
            "         + 0.05(0.9043)",
            "         + 0.08(0.754)",
            "         + 0.05(0.0)",
            "         + 0.02(0.5)",
            "         - 0.13(0.3)",
            "",
            "       = 0.7465",
            "Displayed score: about 75/100",
        ],
        "Maria still ranks strongly because she has full straight-time headroom, strong Long-Term fit, short travel distance, and meaningful tenure."
    )
    add_footer(slide, "This is the exact style of one-employee calculation that makes the recommendation logic easy to explain.")
    add_notes(
        slide,
        """
Now I will show the same employee in the recommendation engine.

Suppose the call-out is a certified Day shift in U-LT1 on April 17, 2026.
Maria Santos is a CNA, so she is eligible for the certified bucket.
For this example we assume she passes the live eligibility checks: no PTO conflict, no unit exclusion, and no blocking source-unit issue.

Now we calculate her normalized values.
Overtime headroom is 1.0 because her current hours are still below the threshold.
Clinical fit is 1.0 because she is Long-Term staff covering another Long-Term unit, which is treated as an exact typology match.
Proximity is one minus 2.87 divided by 30, which is 0.9043.
Float penalty is 0.3 because she is floating to a different unit but the same typology.
Seniority is her tenure divided by 10, capped at 1.0, so 7.54 years becomes 0.754.
Equity is 0 because she is not per diem.
Willingness is the current neutral baseline of 0.5.

Then we multiply each component by its weight and sum them.
The final result is 0.7465, which displays as about 75 out of 100.

This is the most important part of the demo, because it shows that the engine can be audited and explained one person at a time.
        """,
    )


def slide_7(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "What To Emphasize In Q&A", "Three messages that make the demo land")
    add_card(
        slide,
        Inches(0.86), Inches(1.92), Inches(3.72), Inches(3.7),
        "1. Transparent math",
        [
            "Each recommendation can be recalculated by hand.",
            "Weights are explicit and configurable.",
            "The logic is easier to validate with trainers and customers.",
        ],
        fill_rgb=ACCENT_SOFT,
    )
    add_card(
        slide,
        Inches(4.82), Inches(1.92), Inches(3.72), Inches(3.7),
        "2. Operational realism",
        [
            "The engine works from current staffing data.",
            "True shortages stay visible.",
            "The system avoids hiding gaps with unsafe assignments.",
        ],
        fill_rgb=SUCCESS_SOFT,
    )
    add_card(
        slide,
        Inches(8.78), Inches(1.92), Inches(3.58), Inches(3.7),
        "3. Human control",
        [
            "The system recommends and explains.",
            "The coordinator still makes the final call.",
            "Override reasons can be logged for accountability.",
        ],
        fill_rgb=WARN_SOFT,
    )
    add_footer(slide, "Suggested close: this is an explainable staffing engine, not a black box.")
    add_notes(
        slide,
        """
If there is pushback or detailed questioning, I want to close on three points.

First, the math is transparent.
We can calculate one employee's score by hand, which makes the system easier to trust and easier to refine.

Second, the engine is operationally realistic.
It works from the current database state, and if the staffing situation is genuinely constrained, it shows the gap instead of hiding it.

Third, it keeps human control.
The coordinator still makes the decision.
The system is there to make that decision faster, more consistent, and more explainable.
        """,
    )


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    build()
